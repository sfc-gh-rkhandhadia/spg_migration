#!/usr/bin/env python3
"""
Migration Validation Report Generator
=======================================
Generates a Snowflake-branded PowerPoint deck from live SPG validation data.

Usage:
    python3 generate_migration_report.py \
        --client "YourClientName" \
        --author "Rekha Khandhadia" \
        --spg-host "your-spg-host.snowflakecomputing.app" \
        --spg-password "your_password" \
        --run-numbers "1,2,3"          # optional: default = latest 3

All credentials can also be passed via env vars:
    SPG_HOST, SPG_USER, SPG_PASSWORD, SPG_DATABASE
"""

import os, sys, argparse, datetime
import psycopg2, psycopg2.extras
import pymssql

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# ── Colour palette ────────────────────────────────────────────────────────────
DK1       = RGBColor(0x26, 0x26, 0x26)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DK2       = RGBColor(0x11, 0x56, 0x7F)
SF_BLUE   = RGBColor(0x29, 0xB5, 0xE8)
TEAL      = RGBColor(0x71, 0xD3, 0xDC)
ORANGE    = RGBColor(0xFF, 0x9F, 0x36)
VIOLET    = RGBColor(0x7D, 0x44, 0xCF)
PINK      = RGBColor(0xD4, 0x5B, 0x90)
BODY_GREY = RGBColor(0x5B, 0x5B, 0x5B)
TBL_GREY  = RGBColor(0x71, 0x71, 0x71)
LIGHT_BG  = RGBColor(0xF5, 0xF5, 0xF5)
BORDER    = RGBColor(0xC8, 0xC8, 0xC8)
ERROR_RED = RGBColor(0xA2, 0x00, 0x00)
PASS_GREEN= RGBColor(0x1A, 0x7A, 0x3C)   # dark green for pass indicators

DARK_BG_LAYOUTS = {9, 10, 18, 19, 20, 21, 22, 23, 25, 26, 27, 28}
COVER_LAYOUTS   = {13, 14, 15, 16, 17}

# ── Template ──────────────────────────────────────────────────────────────────
TEMPLATE_SEARCH = [
    os.path.join(os.getcwd(), "templates", "snowflake_template.pptx"),
    os.path.expanduser("~/.snowflake/cortex/skills/pptx/snowflake_template.pptx"),
    os.path.expanduser("~/.cortex/skills/900-999_utilities/945-render-pptx/snowflake_template.pptx"),
]
TEMPLATE = next((p for p in TEMPLATE_SEARCH if os.path.isfile(p)), None)
assert TEMPLATE, "snowflake_template.pptx not found"

# ── Helpers ───────────────────────────────────────────────────────────────────

def set_ph(slide, idx, text):
    ph = slide.placeholders[idx]
    t_pos = (ph.top or 0) / 914400
    clean = text.replace('\n', ' ')
    if t_pos < 0.50 and len(clean) > 50:
        print(f"⚠ TITLE TOO LONG: {len(clean)} chars: \"{clean[:50]}...\"")
    ph.text = text
    ph.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    bodyPr = ph.text_frame._txBody.find(f'{{{ns}}}bodyPr')
    if bodyPr is None:
        bodyPr = etree.SubElement(ph.text_frame._txBody, f'{{{ns}}}bodyPr')
    if t_pos < 0.50:
        bodyPr.set('bIns', '0')
    elif 0.60 < t_pos < 1.20:
        bodyPr.set('tIns', '54864')
    if t_pos < 1.20:
        for para in ph.text_frame.paragraphs:
            pPr = para._p.find(f'{{{ns}}}pPr')
            if pPr is None:
                pPr = etree.SubElement(para._p, f'{{{ns}}}pPr')
                para._p.insert(0, pPr)
            pPr.set('indent', '0'); pPr.set('marL', '0')

def _pad_body_ph(ph):
    t_pos = (ph.top or 0) / 914400
    if t_pos > 1.20:
        ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        bodyPr = ph.text_frame._txBody.find(f'{{{ns}}}bodyPr')
        if bodyPr is None:
            bodyPr = etree.SubElement(ph.text_frame._txBody, f'{{{ns}}}bodyPr')
        bodyPr.set('bIns', '91440')

def set_ph_lines(slide, idx, lines, font_size=None):
    ph = slide.placeholders[idx]
    tf = ph.text_frame; tf.clear(); _pad_body_ph(ph)
    lines = [l for l in lines if l.strip()]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        if font_size: p.font.size = Pt(font_size)

def set_ph_sections(slide, idx, sections, heading_size=None, body_size=None):
    ph = slide.placeholders[idx]
    tf = ph.text_frame; tf.clear(); _pad_body_ph(ph)
    first = True
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    for heading, body_lines in sections:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.level = 0
        if not first:
            pPr = p._p.find(f'{{{ns}}}pPr')
            if pPr is None:
                pPr = etree.SubElement(p._p, f'{{{ns}}}pPr'); p._p.insert(0, pPr)
            spcBef = etree.SubElement(pPr, f'{{{ns}}}spcBef')
            spcPts = etree.SubElement(spcBef, f'{{{ns}}}spcPts'); spcPts.set('val', '1400')
        first = False
        run = p.add_run(); run.text = heading; run.font.bold = True
        run.font.color.rgb = DK2
        if heading_size: run.font.size = Pt(heading_size)
        for line in body_lines:
            bp = tf.add_paragraph(); bp.level = 1; bp.text = line
            if body_size: bp.font.size = Pt(body_size)

def add_shape_text(slide, shape_type, left, top, width, height,
                   text, fill_colour, font_colour,
                   font_size=10, bold=False, alignment=PP_ALIGN.CENTER):
    layout_idx = None
    try:
        layout_idx = list(slide.slide_layout.slide_master.slide_layouts).index(slide.slide_layout)
    except (ValueError, AttributeError):
        pass
    if layout_idx in DARK_BG_LAYOUTS or layout_idx in COVER_LAYOUTS:
        if font_colour == DK1:
            font_colour = WHITE
    shape = slide.shapes.add_shape(
        shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_colour
    shape.line.fill.background()
    if width <= 2.0 and '\n' not in text and ' ' in text:
        text = text.replace(' ', '\n')
    tf = shape.text_frame; tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(4); tf.margin_right = Pt(4)
    tf.margin_top = Pt(2);  tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]; p.text = text
    p.font.name = "Arial"; p.font.size = Pt(font_size)
    p.font.bold = bold; p.font.color.rgb = font_colour
    p.alignment = alignment
    return shape

def set_table_borders(tbl, n_rows, n_cols):
    for ri in range(n_rows):
        for ci in range(n_cols):
            tc = tbl.cell(ri, ci)._tc
            tcPr = tc.find(qn("a:tcPr"))
            if tcPr is None: tcPr = etree.SubElement(tc, qn("a:tcPr"))
            for edge in ["lnL", "lnR", "lnT", "lnB"]:
                ln = etree.SubElement(tcPr, qn(f"a:{edge}"), w="12700")
                sf = etree.SubElement(ln, qn("a:solidFill"))
                etree.SubElement(sf, qn("a:srgbClr"), val="C8C8C8")

def add_table_style2(slide, left, top, width, height, headers, data_rows,
                     col_widths=None, font_size=9):
    """Table with DK2 header, alternating rows, BORDER borders."""
    n_rows = len(data_rows) + 1
    n_cols = len(headers)
    tbl_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(left), Inches(top), Inches(width), Inches(height))
    tbl = tbl_shape.table
    if col_widths:
        for ci, cw in enumerate(col_widths):
            tbl.columns[ci].width = Inches(cw)
    # Header
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci); cell.text = h
        cell.fill.solid(); cell.fill.fore_color.rgb = DK2
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(font_size + 1); p.font.bold = True
            p.font.color.rgb = WHITE; p.font.name = "Arial"
            p.alignment = PP_ALIGN.CENTER
    # Data rows
    for ri, row in enumerate(data_rows):
        bg = WHITE if ri % 2 == 0 else LIGHT_BG
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.text = str(val) if val is not None else ''
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size); p.font.name = "Arial"
                # Colour-code verdict column
                txt = str(val)
                if txt in ('PASS','✓ PASS'):
                    p.font.color.rgb = PASS_GREEN; p.font.bold = True
                elif txt in ('FAIL','BOTH_FAILED','SPG_ERROR','SPG_NO_RESULTSET','MSSQL_ONLY'):
                    p.font.color.rgb = ERROR_RED; p.font.bold = True
                else:
                    p.font.color.rgb = TBL_GREY
                p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
    set_table_borders(tbl, n_rows, n_cols)
    return tbl_shape

def verify_slide(slide, prs, slide_num):
    issues = []
    safe_bottom = 5.10
    for shape in slide.shapes:
        l = (shape.left or 0)/914400; t = (shape.top or 0)/914400
        w = (shape.width or 0)/914400; h = (shape.height or 0)/914400
        if not shape.is_placeholder and t+h > safe_bottom and w > 0.5:
            issues.append(f"  OVERFLOW: shape bottom={t+h:.2f}\" > 5.10\"")
        if not shape.is_placeholder and t < 1.22 and w > 0.3 and h > 0.1:
            issues.append(f"  HEADER OVERLAP: shape at ({l:.2f}\",{t:.2f}\") < 1.22\"")
    if issues:
        print(f"⚠ SLIDE {slide_num}: " + " | ".join(issues))
    else:
        print(f"✓ SLIDE {slide_num} OK")
    return issues

def verify_deck(prs):
    content = 0; visual = 0; max_consec = 0; consec = 0
    for slide in prs.slides:
        has_v = any(not s.is_placeholder and (s.width or 0)/914400 > 0.5 and (s.height or 0)/914400 > 0.3
                    for s in slide.shapes)
        is_special = any("Quote" in slide.slide_layout.name or "Thank" in slide.slide_layout.name
                         for _ in [None])
        if not is_special:
            content += 1
            if has_v: visual += 1; consec = 0
            else: consec += 1; max_consec = max(max_consec, consec)
    pct = visual/content*100 if content else 0
    print(f"{'✓' if pct>=40 else '⚠'} DECK: {len(prs.slides)} slides, {visual}/{content} visual ({pct:.0f}%), max {max_consec} consec text")


# ── Data fetching ─────────────────────────────────────────────────────────────

def fetch_mssql_counts(mssql_host, mssql_port, mssql_user, mssql_password, mssql_db):
    """Query MSSQL for actual object counts by schema + type."""
    try:
        conn = pymssql.connect(server=mssql_host, port=int(mssql_port), user=mssql_user,
                               password=mssql_password, database=mssql_db, timeout=15)
        cur = conn.cursor(as_dict=True)
        cur.execute("""
            SELECT s.name AS schema_name,
                CASE o.type
                    WHEN 'U'  THEN 'TABLE'
                    WHEN 'V'  THEN 'VIEW'
                    WHEN 'P'  THEN 'PROCEDURE'
                    WHEN 'TR' THEN 'TRIGGER'
                    ELSE           'FUNCTION'
                END AS object_type,
                COUNT(*) AS cnt
            FROM sys.objects o
            JOIN sys.schemas s ON o.schema_id = s.schema_id
            WHERE o.type IN ('U','V','P','FN','TF','IF','TR')
              AND s.name NOT IN ('sys','INFORMATION_SCHEMA')
            GROUP BY s.name,
                CASE o.type
                    WHEN 'U'  THEN 'TABLE'
                    WHEN 'V'  THEN 'VIEW'
                    WHEN 'P'  THEN 'PROCEDURE'
                    WHEN 'TR' THEN 'TRIGGER'
                    ELSE           'FUNCTION'
                END
        """)
        counts = {(r['schema_name'].lower(), r['object_type']): r['cnt']
                  for r in cur.fetchall()}
        conn.close()
        return counts
    except Exception as e:
        print(f"⚠ MSSQL count fetch failed: {e} — MSSQL column will show n/a")
        return {}


def fetch_spg_counts(spg_host, spg_user, spg_password, spg_db="postgres"):
    """Query SPG catalogs for actual deployed object counts by schema + type."""
    try:
        conn = psycopg2.connect(host=spg_host, port=5432, user=spg_user,
                                password=spg_password, dbname=spg_db,
                                sslmode='require', connect_timeout=20)
        cur = conn.cursor(cursor_factory=psycopg2.extras.RealDictCursor)
        # Exclude all Snowflake/Postgres infrastructure schemas — keep only
        # business migration schemas (api, dbo, stg, etc.)
        _EXCL_SET = (
            'pg_catalog', 'information_schema', 'public', 'validation',
            'cron', 'err', 'incremental', 'map_type',
            '__lake__internal__nsp__', '__pg_lake_table_writes',
            'lake', 'lake_engine', 'lake_file', 'lake_file_cache',
            'lake_iceberg', 'lake_struct', 'lake_table',
            'snowflake_auth', 'snowflake_cdc', 'snowflake_cdc_logs',
            'extension_base',
        )
        _EXCL_IN   = ','.join(f"'{s}'" for s in _EXCL_SET)
        _EXCL_LIKE = (
            "n.nspname NOT LIKE 'pg_%' "
            "AND n.nspname NOT LIKE 'snowflake%' "
            "AND n.nspname NOT LIKE 'lake%' "
            "AND n.nspname NOT LIKE '__lake%' "
            "AND n.nspname NOT LIKE '__pg_%' "
            "AND n.nspname NOT LIKE 'extension%'"
        )
        cur.execute(f"""
            SELECT n.nspname AS s,
                   CASE c.relkind WHEN 'r' THEN 'TABLE' WHEN 'v' THEN 'VIEW' END AS t,
                   COUNT(*) AS cnt
            FROM pg_class c JOIN pg_namespace n ON c.relnamespace=n.oid
            WHERE c.relkind IN ('r','v')
              AND n.nspname NOT IN ({_EXCL_IN})
              AND {_EXCL_LIKE}
            GROUP BY n.nspname, c.relkind
            UNION ALL
            SELECT n.nspname,
                   CASE p.prokind WHEN 'p' THEN 'PROCEDURE' ELSE 'FUNCTION' END,
                   COUNT(*)
            FROM pg_proc p JOIN pg_namespace n ON p.pronamespace=n.oid
            WHERE n.nspname NOT IN ({_EXCL_IN})
              AND {_EXCL_LIKE}
            GROUP BY n.nspname, CASE p.prokind WHEN 'p' THEN 'PROCEDURE' ELSE 'FUNCTION' END
            UNION ALL
            SELECT n.nspname, 'TRIGGER', COUNT(*)
            FROM pg_trigger t
            JOIN pg_class c ON t.tgrelid = c.oid
            JOIN pg_namespace n ON c.relnamespace = n.oid
            WHERE NOT t.tgisinternal
              AND n.nspname NOT IN ({_EXCL_IN})
              AND {_EXCL_LIKE}
            GROUP BY n.nspname
        """)
        counts = {(r['s'].lower(), r['t']): r['cnt'] for r in cur.fetchall()}
        conn.close()
        return counts
    except Exception as e:
        print(f"⚠ SPG catalog count fetch failed: {e} — SPG column will show n/a")
        return {}


def fetch_validation_data(spg_host, spg_user, spg_password, spg_db="postgres",
                          run_numbers=None):
    conn = psycopg2.connect(host=spg_host, port=5432, user=spg_user,
                            password=spg_password, dbname=spg_db,
                            sslmode='require', connect_timeout=20)
    cur = conn.cursor(cursor_factory=psycopg2.extras.RealDictCursor)

    # Latest runs if not specified — one per validation type to avoid duplicates
    if run_numbers is None:
        cur.execute("SELECT run_number, notes FROM validation.validation_run ORDER BY run_number")
        all_runs = cur.fetchall()
        def _latest_by(kw):
            matches = [r['run_number'] for r in all_runs
                       if kw.lower() in (r.get('notes') or '').lower()]
            return matches[-1] if matches else None
        rn_t = _latest_by('trigger')
        rn_v = _latest_by('view')
        rn_p = _latest_by('output comparison')
        rn_w = _latest_by('rollback-wrapped')
        run_numbers = sorted(set(r for r in [rn_t, rn_v, rn_p, rn_w] if r))

    rn_list = ','.join(str(r) for r in run_numbers)

    # Run summary
    cur.execute(f"""
        SELECT * FROM validation.validation_run
        WHERE run_number IN ({rn_list}) ORDER BY run_number
    """)
    runs = cur.fetchall()

    # Schema summary (with table counts from pg catalogs)
    # Deduplicate: when the same object appears as SKIPPED in one run and tested in another,
    # prefer the non-SKIPPED (actual result) row using ROW_NUMBER ranked non-SKIPPED first.
    cur.execute(f"""
        SELECT source_schema, object_type,
               COUNT(*) AS total,
               COUNT(*) FILTER (WHERE test_verdict IN ('PASS','PASS_DML_PROC','PASS_WRITE_PROC','WRITE_EXPECTED_FAIL')) AS passed,
               COUNT(*) FILTER (WHERE test_verdict NOT IN ('PASS','PASS_DML_PROC','PASS_WRITE_PROC','WRITE_EXPECTED_FAIL','SKIPPED','FAIL_MISSING_PREREQ')) AS failed,
               COUNT(*) FILTER (WHERE test_verdict='SKIPPED') AS skipped
        FROM (
            SELECT *,
                   ROW_NUMBER() OVER (
                       PARTITION BY source_schema, object_name
                       ORDER BY CASE WHEN test_verdict != 'SKIPPED' THEN 0 ELSE 1 END,
                                run_number DESC
                   ) AS _rn
            FROM validation.validation_result
            WHERE run_number IN ({rn_list})
        ) deduped
        WHERE _rn = 1
        GROUP BY source_schema, object_type
        ORDER BY source_schema, object_type
    """)
    schema_summary = cur.fetchall()

    # Failure details by schema
    cur.execute(f"""
        SELECT source_schema, object_type, object_name, test_verdict,
               COALESCE(issues[1], error_message, '') AS top_issue
        FROM validation.validation_result
        WHERE run_number IN ({rn_list})
          AND test_verdict NOT IN ('PASS','SKIPPED','PASS_DML_PROC','PASS_WRITE_PROC','WRITE_EXPECTED_FAIL')
        ORDER BY source_schema, object_type, test_verdict, object_name
    """)
    failures = cur.fetchall()

    # Verdict category counts
    cur.execute(f"""
        SELECT test_verdict, COUNT(*) AS cnt
        FROM validation.validation_result
        WHERE run_number IN ({rn_list})
          AND test_verdict NOT IN ('PASS','SKIPPED','PASS_DML_PROC','PASS_WRITE_PROC','WRITE_EXPECTED_FAIL')
        GROUP BY test_verdict ORDER BY cnt DESC
    """)
    verdict_cats = cur.fetchall()

    conn.close()
    return runs, schema_summary, failures, verdict_cats, run_numbers


# ── KPI aggregation ───────────────────────────────────────────────────────────

def compute_kpis(runs, schema_summary, mssql_counts, spg_catalog_counts):
    # Use the deduplicated schema_summary (already de-duped by object name across runs)
    # rather than summing per-run pass_count/fail_count which would double-count objects
    # that were SKIPPED in one run and tested in another.
    total_pass = sum((r.get('passed') or 0) for r in schema_summary
                     if r.get('source_schema') not in ('validation',))
    total_fail = sum((r.get('failed') or 0) for r in schema_summary
                     if r.get('source_schema') not in ('validation',))
    tested = total_pass + total_fail
    pass_rate = round(total_pass / tested * 100) if tested else 0
    schemas = sorted({r['source_schema'] for r in schema_summary
                      if r['source_schema'] not in ('validation',)})
    # Use live catalog counts for totals (excludes system schemas)
    total_mssql = sum(v for (s, t), v in mssql_counts.items()
                      if s not in ('sys', 'information_schema'))
    total_spg   = sum(v for (s, t), v in spg_catalog_counts.items()
                      if s not in ('validation',))
    return {
        'total_mssql_objects': total_mssql or sum(r.get('total_objects') or 0 for r in runs),
        'spg_objects': total_spg or sum(r['total'] for r in schema_summary),
        'pass_rate': pass_rate,
        'schemas': len(schemas),
        'schemas_list': schemas,
        'total_pass': total_pass,
        'total_fail': total_fail,
    }


# ── Slide builders ────────────────────────────────────────────────────────────

def slide_cover(prs, client_name, run_date, author):
    slide = prs.slides.add_slide(prs.slide_layouts[13])
    set_ph(slide, 3, "MSSQL → SNOWFLAKE\nPOSTGRES MIGRATION")
    set_ph(slide, 0, f"Migration Validation Report  |  {client_name}")
    set_ph(slide, 2, f"{author}  |  {run_date}")
    verify_slide(slide, prs, 1)


def slide_testing_approach(prs, client_name="", run_date="", mssql_db=""):
    """Slide 2 — Executive summary of the 4-step migration testing approach."""
    slide   = prs.slides.add_slide(prs.slide_layouts[0])
    set_ph(slide, 0, "TESTING APPROACH")
    subtitle = f"How the {mssql_db} MSSQL → Snowflake Postgres migration was validated"
    if client_name:
        subtitle += f"  |  Client: {client_name}"
    if run_date:
        subtitle += f"  |  {run_date}"
    set_ph(slide, 1, subtitle)

    slide_w = prs.slide_width.inches   # 13.33" widescreen

    steps = [
        ("1", "Source Database Setup",
         "Created the source MSSQL database in a local Docker environment "
         "using the client\u2019s own source code and schema scripts."),
        ("2", "Random Data Loading",
         "Loaded the database with representative random data using "
         "Cortex Code (CoCo) to simulate production volume for meaningful comparison."),
        ("3", "Code Migration",
         "Used the Cortex Code MSSQL \u2192 Snowflake Postgres Migration Skill "
         "to convert all tables, views, procedures, functions and triggers."),
        ("4", "Migration Validation",
         "Used the Cortex Code Validation Skill to execute both systems "
         "side-by-side and compare outputs, row counts, schemas and behaviour."),
    ]

    # 2 columns × 2 rows of step cards
    card_gap_x = 0.45
    card_gap_y = 0.30
    n_cols     = 2
    card_w     = (slide_w - 1.20 - (n_cols - 1) * card_gap_x) / n_cols  # ≈ 5.64"
    card_h     = 1.62
    start_x    = (slide_w - n_cols * card_w - (n_cols - 1) * card_gap_x) / 2.0
    start_y    = 1.45

    for i, (num, heading, body) in enumerate(steps):
        col   = i % n_cols
        row   = i // n_cols
        cx    = start_x + col * (card_w + card_gap_x)
        cy    = start_y + row * (card_h  + card_gap_y)

        # Card background
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(cx), Inches(cy), Inches(card_w), Inches(card_h))
        card.fill.solid()
        card.fill.fore_color.rgb = DK2
        card.line.fill.background()
        card.adjustments[0] = 0.04   # corner radius

        # Step number badge (top-left)
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(cx + 0.18), Inches(cy + 0.14), Inches(0.46), Inches(0.46))
        badge.fill.solid(); badge.fill.fore_color.rgb = SF_BLUE
        badge.line.fill.background()
        badge.adjustments[0] = 0.25
        btf = badge.text_frame
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        bp = btf.paragraphs[0]; bp.text = num
        bp.font.name = "Arial"; bp.font.size = Pt(16); bp.font.bold = True
        bp.font.color.rgb = WHITE; bp.alignment = PP_ALIGN.CENTER

        # Heading
        htb = slide.shapes.add_textbox(
            Inches(cx + 0.76), Inches(cy + 0.14), Inches(card_w - 0.92), Inches(0.38))
        htf = htb.text_frame; htf.word_wrap = False
        hp = htf.paragraphs[0]; hp.text = heading
        hp.font.name = "Arial"; hp.font.size = Pt(13); hp.font.bold = True
        hp.font.color.rgb = SF_BLUE; hp.alignment = PP_ALIGN.LEFT
        htf.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Divider line under heading
        div = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(cx + 0.18), Inches(cy + 0.58), Inches(card_w - 0.36), Inches(0.02))
        div.fill.solid(); div.fill.fore_color.rgb = SF_BLUE
        div.line.fill.background()

        # Body text
        btb = slide.shapes.add_textbox(
            Inches(cx + 0.18), Inches(cy + 0.66), Inches(card_w - 0.36), Inches(card_h - 0.80))
        btf2 = btb.text_frame; btf2.word_wrap = True
        btp = btf2.paragraphs[0]; btp.text = body
        btp.font.name = "Arial"; btp.font.size = Pt(11); btp.font.bold = False
        btp.font.color.rgb = WHITE; btp.alignment = PP_ALIGN.LEFT
        btf2.vertical_anchor = MSO_ANCHOR.TOP

    verify_slide(slide, prs, 2)


def slide_methodology(prs):
    ns   = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    slide_w = prs.slide_width.inches
    COL_W   = min(9.00, slide_w - 1.20)   # full-width content column
    COL_TOP = 1.30
    COL_H   = 3.75
    col_x   = (slide_w - COL_W) / 2.0     # centred

    def zero_spacing(p, spc_aft_pt=0):
        pPr = p._p.find(f"{{{ns}}}pPr")
        if pPr is None:
            pPr = etree.SubElement(p._p, f"{{{ns}}}pPr")
            p._p.insert(0, pPr)
        for tag in ("spcBef", "spcAft", "lnSpc"):
            for old in pPr.findall(f"{{{ns}}}{tag}"):
                pPr.remove(old)
        b = etree.SubElement(pPr, f"{{{ns}}}spcBef")
        etree.SubElement(b, f"{{{ns}}}spcPts").set("val", "0")
        a = etree.SubElement(pPr, f"{{{ns}}}spcAft")
        etree.SubElement(a, f"{{{ns}}}spcPts").set("val", str(spc_aft_pt * 100))
        ln = etree.SubElement(pPr, f"{{{ns}}}lnSpc")
        etree.SubElement(ln, f"{{{ns}}}spcPct").set("val", "115000")

    def build_slide(title, subtitle, heading, bullets, slide_num):
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        set_ph(slide, 0, title)
        set_ph(slide, 1, subtitle)

        shape = slide.shapes.add_textbox(
            Inches(col_x), Inches(COL_TOP), Inches(COL_W), Inches(COL_H))
        tf = shape.text_frame
        tf.word_wrap     = True
        tf.auto_size     = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.margin_top    = Inches(0)
        tf.margin_bottom = Inches(0)
        tf.margin_left   = Inches(0.10)
        tf.margin_right  = Inches(0)

        # Section heading
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = heading
        r.font.bold = True; r.font.color.rgb = DK2
        r.font.name = "Arial"; r.font.size = Pt(18)
        p.alignment = PP_ALIGN.CENTER
        zero_spacing(p, spc_aft_pt=8)

        # Underline bar
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(col_x), Inches(COL_TOP + 0.30),
            Inches(COL_W), Inches(0.03))
        bar.fill.solid(); bar.fill.fore_color.rgb = DK2
        bar.line.fill.background()

        # Bullets
        for line in bullets:
            bp = tf.add_paragraph()
            bp.alignment = PP_ALIGN.LEFT
            zero_spacing(bp, spc_aft_pt=5)
            pPr = bp._p.find(f"{{{ns}}}pPr")
            pPr.set("marL", "228600")
            pPr.set("indent", "-228600")
            r2 = bp.add_run()
            r2.text = "\u2022  " + line
            r2.font.name = "Arial"; r2.font.size = Pt(14)
            r2.font.color.rgb = DK1

        verify_slide(slide, prs, slide_num)

    # ── Slide 2: Objects Validated ────────────────────────────────────────────
    build_slide(
        title    = "OBJECTS VALIDATED",
        subtitle = "What was tested across both source and target databases",
        heading  = "Objects Validated",
        bullets  = [
            "Views — schema, row count and data hash",
            "Stored Procedures — parameter count and execution parity",
            "Functions — existence, parameter parity and output",
            "Triggers — events, table target and enabled state",
            "Tables — MSSQL vs SPG object count",
            "Constraints — PRIMARY KEY, FOREIGN KEY, UNIQUE, CHECK",
            "Results stored for validation",
            "Summarize validation results",
        ],
        slide_num = 4,
    )

    # ── Slide 3: Validation Approach ──────────────────────────────────────────
    build_slide(
        title    = "VALIDATION APPROACH",
        subtitle = "How each object type was executed and compared",
        heading  = "Validation Approach",
        bullets  = [
            "MSSQL and SPG executed with identical real sampled parameters",
            "Row count, column names and data hash compared side-by-side",
            "Verdicts: PASS / FAIL / SPG_ERROR / BOTH_FAILED / SKIPPED",
            "Write/modify procedures excluded from execution (safe mode)",
            "Views: up to 2,000 rows sampled and MD5-hashed per side",
            "Column schema verified per view — missing columns flagged",
            "ROW_COUNT_MISMATCH and DATA_HASH_MISMATCH classified separately",
        ],
        slide_num = 3,
    )


def slide_section_divider(prs, part_label, title, subtitle=""):
    """Dark full-bleed section divider slide (Part 1 / Part 2 separator)."""
    slide = prs.slides.add_slide(prs.slide_layouts[18])  # dark/chapter layout
    set_ph(slide, 1, f"{part_label}\n{title}")
    if subtitle:
        tb = slide.shapes.add_textbox(
            Inches(0.60), Inches(3.80), Inches(12.10), Inches(0.55))
        tp = tb.text_frame.paragraphs[0]
        tp.text = subtitle
        tp.font.name  = "Arial"
        tp.font.size  = Pt(14)
        tp.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        tp.alignment  = PP_ALIGN.CENTER
    verify_slide(slide, prs, 99)


# ─── Part 1 helpers ───────────────────────────────────────────────────────────

def slide_part1_schema_deployment(prs, mssql_counts, spg_catalog_counts):
    """Part 1 Slide: structural deployment summary — Object Type | MSSQL | SPG | Deployed | Missing | Extra | Status."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_ph(slide, 0, "PART 1 — SCHEMA DEPLOYMENT SUMMARY")
    set_ph(slide, 1, "Checks whether each MSSQL object was deployed to Snowflake Postgres")

    OBJ_TYPES = ['TABLE', 'VIEW', 'PROCEDURE', 'FUNCTION', 'TRIGGER', 'INDEX']
    EXCL_SCH  = {'sys', 'information_schema', 'validation'}
    BUSINESS  = {'api', 'stg', 'dbo', 'err', 'svc_menu_management'}

    rows = []
    for otype in OBJ_TYPES:
        ms_cnt  = sum(v for (s, t), v in mssql_counts.items()
                      if t == otype and s not in EXCL_SCH)
        spg_cnt = sum(v for (s, t), v in spg_catalog_counts.items()
                      if t == otype and s not in EXCL_SCH and s in BUSINESS)
        deployed = min(ms_cnt, spg_cnt)
        missing  = max(ms_cnt - spg_cnt, 0)
        extra    = max(spg_cnt - ms_cnt, 0)
        if missing == 0:
            status = u"\u2705 All deployed" + (f" (+{extra} extra)" if extra else "")
        else:
            status = f"\u26a0\ufe0f {missing} not deployed"
        rows.append([
            otype,
            str(ms_cnt) or "0",
            str(spg_cnt) or "0",
            str(deployed) or "0",
            str(missing) or "0",
            str(extra) or "0",
            status,
        ])

    headers    = ["Object Type", "# MSSQL", "# SPG", "Deployed", "Missing", "Extra in SPG", "Status"]
    col_widths = [1.50, 0.90, 0.80, 0.90, 0.90, 1.10, 2.00]
    tbl_w      = sum(col_widths)
    tbl_x      = (prs.slide_width.inches - tbl_w) / 2.0
    add_table_style2(slide, tbl_x, 1.35, tbl_w, min(3.50, 0.45 * (len(rows)+1)),
                     headers, rows, col_widths=col_widths, font_size=12)

    fn = slide.shapes.add_textbox(Inches(tbl_x), Inches(4.70), Inches(tbl_w), Inches(0.28))
    fn.text_frame.paragraphs[0].text = (
        u"\u24d8  Extra = additional objects added by the converter beyond the MSSQL source — not a gap."
    )
    fn.text_frame.paragraphs[0].font.size = Pt(7)
    fn.text_frame.paragraphs[0].font.italic = True
    fn.text_frame.paragraphs[0].font.color.rgb = BODY_GREY
    verify_slide(slide, prs, 100)


def slide_part1_by_schema(prs, mssql_counts, spg_catalog_counts):
    """Part 1 Slide: structural coverage per business schema."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_ph(slide, 0, "PART 1 — STRUCTURAL COVERAGE BY SCHEMA")
    set_ph(slide, 1, "deployed / total per schema (\u2705 all deployed  \u26a0 gaps)")

    OBJ_TYPES = ['TABLE', 'VIEW', 'PROCEDURE', 'FUNCTION', 'TRIGGER', 'INDEX']
    BUSINESS  = sorted({'api', 'stg', 'dbo', 'err'} &
                       {s for (s, t) in mssql_counts})
    EXCL_SCH  = {'sys', 'information_schema'}

    rows = []
    for sc in BUSINESS:
        row = [f"`{sc}`"]
        all_ok = True
        for otype in OBJ_TYPES:
            ms  = mssql_counts.get((sc, otype), 0)
            spg = spg_catalog_counts.get((sc, otype), 0)
            if ms == 0:
                cell = u"\u2014"
            else:
                dep = min(ms, spg)
                if dep < ms:
                    all_ok = False
                cell = f"{dep}/{ms}"
            row.append(cell)
        row.append(u"\u2705" if all_ok else u"\u26a0")
        rows.append(row)

    headers    = ["Schema"] + OBJ_TYPES + ["Status"]
    col_widths = [0.90] + [1.05] * len(OBJ_TYPES) + [0.70]
    tbl_w      = sum(col_widths)
    tbl_x      = (prs.slide_width.inches - tbl_w) / 2.0
    tbl_shape  = add_table_style2(slide, tbl_x, 1.35, tbl_w, min(3.20, 0.45 * (len(rows)+1)),
                     headers, rows, col_widths=col_widths, font_size=12)

    # Post-process: colour ⚠ status cells yellow, bold and larger
    WARN_YELLOW = RGBColor(0xFF, 0xC0, 0x00)
    tbl = tbl_shape.table
    status_col = len(headers) - 1
    for ri in range(1, len(rows) + 1):
        cell = tbl.cell(ri, status_col)
        txt  = cell.text.strip()
        if '\u26a0' in txt:            # ⚠ warning icon
            cell.fill.solid()
            cell.fill.fore_color.rgb = WARN_YELLOW
            for p in cell.text_frame.paragraphs:
                p.font.size  = Pt(18)
                p.font.bold  = True
                p.font.color.rgb = RGBColor(0x33, 0x33, 0x00)  # dark on yellow
        elif '\u2705' in txt:          # ✅ pass icon — keep green tint
            for p in cell.text_frame.paragraphs:
                p.font.size  = Pt(16)
                p.font.bold  = True

    verify_slide(slide, prs, 100)


def slide_exec_summary(prs, kpis):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_ph(slide, 0, "PART 2 — BEHAVIORAL VALIDATION")
    set_ph(slide, 1, f"Live execution parity across {kpis['schemas']} schemas — triggers, views, procedures, functions")

    stats = [
        (str(kpis['total_mssql_objects']), "TOTAL MSSQL\nOBJECTS", DK2),
        (str(kpis['spg_objects']),         "DEPLOYED TO\nSNOWFLAKE PG", SF_BLUE),
        (f"{kpis['pass_rate']}%",          "OVERALL\nPASS RATE", PASS_GREEN if kpis['pass_rate'] >= 70 else ORANGE),
        (str(kpis['schemas']),             "SCHEMAS\nTESTED", VIOLET),
    ]
    box_w = 2.00; gap = 0.22; x = 0.40
    for val, label, colour in stats:
        # Big number
        num_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(1.45), Inches(box_w), Inches(1.30))
        num_box.fill.solid(); num_box.fill.fore_color.rgb = colour
        num_box.line.fill.background()
        tf = num_box.text_frame; tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.text = val
        p.font.name = "Arial"; p.font.size = Pt(32); p.font.bold = True
        p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
        # Label below
        lbl = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(2.80), Inches(box_w), Inches(0.55))
        lbl.fill.solid(); lbl.fill.fore_color.rgb = DK2
        lbl.line.fill.background()
        tf2 = lbl.text_frame; tf2.word_wrap = True
        tf2.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
        p2 = tf2.paragraphs[0]; p2.text = label
        p2.font.name = "Arial"; p2.font.size = Pt(9); p2.font.bold = True
        p2.font.color.rgb = WHITE; p2.alignment = PP_ALIGN.CENTER
        x += box_w + gap

    # Pass/Fail sub-line
    sub = slide.shapes.add_textbox(Inches(0.40), Inches(3.50), Inches(9.10), Inches(0.30))
    p = sub.text_frame.paragraphs[0]
    p.text = f"Objects Passed: {kpis['total_pass']}   |   Objects Failed: {kpis['total_fail']}   |   Schemas: {', '.join(kpis['schemas_list'])}"
    p.font.name = "Arial"; p.font.size = Pt(10); p.font.color.rgb = BODY_GREY
    p.alignment = PP_ALIGN.CENTER

    verify_slide(slide, prs, 5)


def slide_schema_table(prs, schema_summary, mssql_counts, spg_catalog_counts):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_ph(slide, 0, "BEHAVIORAL RESULTS BY SCHEMA + TYPE")
    set_ph(slide, 1, "Execution parity — triggers, views, procedures and functions")

    headers = ["Schema", "Type", "# MSSQL", "# SPG", "Passed", "Failed", "Pass %"]
    # Collect all (schema, type) keys from all three sources
    _BEHAVIORAL = {'PROCEDURE', 'FUNCTION', 'VIEW', 'TRIGGER'}
    all_keys = sorted(
        {(r['source_schema'], r['object_type']) for r in schema_summary
         if r['source_schema'] not in ('validation',)
         and r['object_type'] in _BEHAVIORAL}
        | {(s, t) for (s, t) in mssql_counts
           if s not in ('sys','information_schema') and t in _BEHAVIORAL}
        | {(s, t) for (s, t) in spg_catalog_counts
           if s not in ('validation',) and t in _BEHAVIORAL}
    )
    data_rows = []
    for schema, otype in all_keys:
        # MSSQL count from live query
        ms_cnt = mssql_counts.get((schema, otype), 0)
        # SPG deployed count from catalog
        spg_cnt = spg_catalog_counts.get((schema, otype), 0)
        # Validation results (pass/fail) — may differ from catalog if objects weren't testable
        val_row = next((r for r in schema_summary
                        if r['source_schema'] == schema and r['object_type'] == otype), None)
        if val_row:
            # Execution-validated row (PROCEDURE/FUNCTION/TRIGGER)
            passed = val_row['passed'] or 0
            failed = val_row['failed'] or 0
            tested = passed + failed
            pct = f"{round(passed/tested*100)}%" if tested > 0 else "n/a"
        else:
            # TABLE/VIEW: not in validation_result — use structural (catalog) match
            passed = min(ms_cnt, spg_cnt)
            failed = max(ms_cnt - spg_cnt, 0)
            pct = f"{round(passed/ms_cnt*100)}%" if ms_cnt > 0 else "n/a"
        data_rows.append([
            schema, otype,
            str(ms_cnt) if ms_cnt else "0",
            str(spg_cnt) if spg_cnt else "0",
            str(passed), str(failed), pct
        ])

    # Fit all rows onto a single slide — no pagination
    col_widths = [1.30, 1.70, 0.85, 0.85, 0.85, 0.85, 0.85]  # wider Type col, total 8.25"
    tbl_w  = sum(col_widths)                              # 8.25"
    tbl_x  = (prs.slide_width.inches - tbl_w) / 2.0      # centred on slide
    n_rows = len(data_rows) + 1
    row_h  = min(0.29, 2.70 / max(n_rows, 1))            # keep all rows above footnote
    tbl_h  = n_rows * row_h
    add_table_style2(slide, tbl_x, 1.55, tbl_w, tbl_h, headers, data_rows,
                     col_widths=col_widths, font_size=11)
    # Footnote — fixed at bottom, never inside the table
    fn_y = max(1.55 + tbl_h + 0.18, 4.70)
    fn_y = min(fn_y, 4.80)
    fn = slide.shapes.add_textbox(Inches(tbl_x), Inches(fn_y), Inches(tbl_w), Inches(0.35))
    fn.text_frame.word_wrap = True
    fp = fn.text_frame.paragraphs[0]
    fp.text = (
        u"\u24d8  Passed = PASS + PASS_DML_PROC + PASS_WRITE_PROC + WRITE_EXPECTED_FAIL. "
        "SPG shows more FUNCTIONs than MSSQL because MSSQL reader procedures that return rows "
        "are correctly migrated as PostgreSQL FUNCTIONS. This is expected and not a defect."
    )
    fp.font.name = "Arial"; fp.font.size = Pt(7)
    fp.font.color.rgb = BODY_GREY; fp.font.italic = True
    fp.alignment = PP_ALIGN.LEFT
    verify_slide(slide, prs, 6)


def slide_pass_rate_visual(prs, schema_summary, mssql_counts, spg_catalog_counts):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_ph(slide, 0, "PASS RATE BY SCHEMA")
    set_ph(slide, 1, "Aggregated pass rate across all tested object types per schema (TABLEs excluded)")

    # Aggregate pass/fail by schema — skip TABLE rows (not validated) and infra schemas
    schema_totals = {}
    for r in schema_summary:
        sc = r['source_schema']
        if sc in ('validation',): continue
        if r['object_type'] == 'TABLE': continue
        if sc not in schema_totals:
            schema_totals[sc] = {'passed': 0, 'failed': 0}
        schema_totals[sc]['passed'] += (r['passed'] or 0)
        schema_totals[sc]['failed'] += (r['failed'] or 0)

    rows = [{'source_schema': sc, 'passed': v['passed'], 'failed': v['failed']}
            for sc, v in sorted(schema_totals.items())]

    n = len(rows)
    # Reserve bottom 1.0" for legend — bars share the upper area
    bar_area_h = 2.80
    bar_h = min(0.55, bar_area_h / max(n, 1))
    gap   = 0.22
    y     = 1.90
    label_w   = 1.80   # schema name label
    bar_max_w = 7.00   # progress bar width
    count_w   = 0.80   # pct label
    pct_gap   = 0.12   # gap between bar end and pct label
    # Centre using actual slide width (widescreen = 13.33", standard = 10")
    slide_w   = prs.slide_width.inches
    total_w   = label_w + 0.04 + bar_max_w + pct_gap + count_w
    margin_x  = (slide_w - total_w) / 2.0

    for r in rows:
        passed = r['passed'] or 0
        failed = r['failed'] or 0
        tested = passed + failed
        pct    = passed / tested if tested > 0 else 0

        # Schema label
        lbl = slide.shapes.add_textbox(
            Inches(margin_x), Inches(y), Inches(label_w), Inches(bar_h))
        lbl.text_frame.word_wrap = False
        p = lbl.text_frame.paragraphs[0]
        p.text = r['source_schema']
        p.font.name = "Arial"; p.font.size = Pt(14); p.font.bold = True
        p.font.color.rgb = DK2; p.alignment = PP_ALIGN.CENTER
        lbl.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        bar_x = margin_x + label_w + 0.04

        # Background bar (light grey = total / not validated)
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(bar_x), Inches(y + 0.05), Inches(bar_max_w), Inches(bar_h - 0.10))
        bg.fill.solid(); bg.fill.fore_color.rgb = LIGHT_BG
        bg.line.fill.background()

        # Pass bar (blue)
        if pct > 0:
            pw = max(0.05, bar_max_w * pct)
            pb = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(bar_x), Inches(y + 0.05), Inches(pw), Inches(bar_h - 0.10))
            pb.fill.solid(); pb.fill.fore_color.rgb = SF_BLUE
            pb.line.fill.background()

        # Fail bar (orange, right portion)
        if pct < 1.0 and failed > 0:
            fw = bar_max_w * (1 - pct)
            fb = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(bar_x + bar_max_w * pct), Inches(y + 0.05),
                Inches(fw), Inches(bar_h - 0.10))
            fb.fill.solid(); fb.fill.fore_color.rgb = ORANGE
            fb.line.fill.background()

        # Pass counts inside bar
        inside = slide.shapes.add_textbox(
            Inches(bar_x + 0.05), Inches(y + 0.05), Inches(bar_max_w - 0.10), Inches(bar_h - 0.10))
        inside.text_frame.word_wrap = False
        ip = inside.text_frame.paragraphs[0]
        ip.text = f"{passed} passed  |  {failed} failed  |  {tested} tested"
        ip.font.name = "Arial"; ip.font.size = Pt(11); ip.font.bold = True
        ip.font.color.rgb = WHITE if pct > 0.3 else DK2
        ip.alignment = PP_ALIGN.CENTER
        inside.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Percentage label — right of bar
        pct_lbl = slide.shapes.add_textbox(
            Inches(bar_x + bar_max_w + pct_gap),
            Inches(y), Inches(count_w), Inches(bar_h))
        pct_lbl.text_frame.word_wrap = False
        p2 = pct_lbl.text_frame.paragraphs[0]
        pct_str = f"{round(pct*100)}%" if tested > 0 else "n/a"
        p2.text = pct_str
        p2.font.name = "Arial"; p2.font.size = Pt(16); p2.font.bold = True
        p2.font.color.rgb = (SF_BLUE if pct >= 0.9
                             else ORANGE if pct >= 0.5
                             else ERROR_RED)
        p2.alignment = PP_ALIGN.LEFT
        pct_lbl.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        y += bar_h + gap

    # ── Legend — fixed at bottom, centred on slide ───────────────────────────
    legend_y   = 4.62
    swatch_w   = 0.22; swatch_h = 0.22; item_gap = 0.18; label_gap = 0.08
    text_w     = 1.50
    _item_w    = swatch_w + label_gap + text_w          # width of one item
    _n_items   = 3
    _legend_w  = _n_items * _item_w + (_n_items - 1) * item_gap  # total legend width
    legend_x   = (slide_w - _legend_w) / 2.0            # centred on slide

    for colour, text in [(SF_BLUE, "PASSED"), (ORANGE, "FAILED / ERROR"),
                         (LIGHT_BG, "NOT VALIDATED")]:
        # Colour swatch
        sw = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(legend_x), Inches(legend_y), Inches(swatch_w), Inches(swatch_h))
        sw.fill.solid(); sw.fill.fore_color.rgb = colour
        sw.line.fill.background()
        # Text label beside swatch
        txt = slide.shapes.add_textbox(
            Inches(legend_x + swatch_w + label_gap), Inches(legend_y - 0.02),
            Inches(text_w), Inches(swatch_h + 0.04))
        tp = txt.text_frame.paragraphs[0]
        tp.text = text
        tp.font.name = "Arial"; tp.font.size = Pt(10); tp.font.bold = True
        tp.font.color.rgb = DK2
        legend_x += _item_w + item_gap

    verify_slide(slide, prs, 9)


def slide_failure_categories(prs, verdict_cats):
    slide = prs.slides.add_slide(prs.slide_layouts[7])
    set_ph(slide, 0, "FAILURE CATEGORIES")
    set_ph(slide, 4, "Breakdown of validation issues requiring remediation")

    # Group into 3 columns
    missing  = [(r['test_verdict'], r['cnt']) for r in verdict_cats if r['test_verdict'] == 'MSSQL_ONLY']
    spg_err  = [(r['test_verdict'], r['cnt']) for r in verdict_cats if r['test_verdict'] in ('SPG_ERROR','SPG_NO_RESULTSET','FAIL')]
    both_fail= [(r['test_verdict'], r['cnt']) for r in verdict_cats if r['test_verdict'] == 'BOTH_FAILED']

    miss_total = sum(c for _, c in missing)
    spg_total  = sum(c for _, c in spg_err)
    both_total = sum(c for _, c in both_fail)

    set_ph_sections(slide, 1, [
        ("Missing in SPG", [f"{miss_total} objects"]),
        ("Root Cause", [
            "Not yet deployed to Snowflake PG",
            "Schema conversion incomplete",
            "Deploy from migration package",
        ]),
    ], heading_size=11, body_size=10)

    set_ph_sections(slide, 2, [
        ("SPG Execution Errors", [f"{spg_total} objects"]),
        ("Root Cause", [
            "boolean = integer comparison",
            "PROCEDURE vs FUNCTION type mismatch",
            "Column name typos in migration",
            "Parameter signature differences",
        ]),
    ], heading_size=11, body_size=10)

    set_ph_sections(slide, 3, [
        ("Both Sides Failed", [f"{both_total} objects"]),
        ("Root Cause", [
            "MSSQL requires active job/data",
            "SPG parameter name mismatch",
            "stg.p_microsloadlog_update",
            "Fix param names → unblocks 34",
        ]),
    ], heading_size=11, body_size=10)

    verify_slide(slide, prs, 8)


def slide_failure_detail(prs, failures, schema_filter, slide_num, title_suffix=""):
    filtered = [f for f in failures if f['source_schema'] in schema_filter]
    if not filtered:
        return slide_num

    headers = ["Object Name", "Type", "Verdict", "Issue"]
    page_size = 7
    for page_idx in range(0, max(1, len(filtered)), page_size):
        page = filtered[page_idx:page_idx+page_size]
        sl = prs.slides.add_slide(prs.slide_layouts[0])
        schemas_str = ' + '.join(f'`{s}`' for s in schema_filter)
        suffix = f" (cont.)" if page_idx > 0 else ""
        set_ph(sl, 0, f"FAILED OBJECTS — {', '.join(s.upper() for s in schema_filter)}{title_suffix}")
        set_ph(sl, 1, f"Objects requiring remediation{suffix} — {len(filtered)} total")
        data_rows = []
        for f in page:
            obj = f['object_name'] or ''
            if '.' in obj: obj = obj.split('.')[-1]
            issue = (f['top_issue'] or '')[:60]
            data_rows.append([obj[:35], f['object_type'], f['test_verdict'], issue])
        n_rows = len(data_rows) + 1
        row_h = min(0.40, 3.60 / n_rows)
        add_table_style2(sl, 0.40, 1.35, 9.10, n_rows * row_h, headers, data_rows,
                         col_widths=[2.40, 1.10, 1.40, 4.20], font_size=9)
        verify_slide(sl, prs, slide_num)
        slide_num += 1
    return slide_num


def slide_remediation(prs, verdict_cats):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_ph(slide, 0, "REMEDIATION PRIORITIES")
    set_ph(slide, 1, "Ranked by number of objects unblocked per fix")

    steps = [
        ("1", "Fix stg.p_microsloadlog_update\nparam names (v_ → p_)", DK2,    "Unblocks 34\nBOTH_FAILED"),
        ("2", "Convert 6 api PROCEDUREs\nto FUNCTION RETURNS TABLE", SF_BLUE,  "Fixes\nSPG_NO_RESULTSET"),
        ("3", "Fix boolean = integer\n(iscpg, isactive, jobactive)",   DK2,    "Fixes 20+\nSPG_ERROR"),
        ("4", "Fix dbo.p_errorlog\nsignature mismatch",                VIOLET, "Unblocks\n4 dbo objects"),
        ("5", "Deploy 23 missing\nprocedures to SPG",                  SF_BLUE,"Removes\nMSSQL_ONLY"),
    ]

    # Layout: 5 boxes + 4 arrows, centred horizontally and vertically on widescreen slide (13.33" × 7.5")
    slide_w  = prs.slide_width.inches
    slide_h  = prs.slide_height.inches
    step_w   = 1.50   # narrower to fit all 5 boxes horizontally with safe margins
    step_h   = 1.15   # compact height to keep impact labels within slide
    arrow_w  = 0.30
    arrow_h  = 0.30
    arrow_gap= 0.08
    badge_h  = 0.40

    n        = len(steps)
    total_w  = n * step_w + (n - 1) * (arrow_gap * 2 + arrow_w)
    x_start  = (slide_w - total_w) / 2.0   # ≈1.26" margin each side

    # Vertically centre the content block on the full slide height
    content_h = badge_h + 0.30 + step_h + 0.22 + 0.55
    y_box     = (slide_h - content_h) / 2.0   # true visual centre
    y_label   = y_box + badge_h + 0.30 + step_h + 0.22

    x = x_start
    for num, action, colour, impact in steps:
        badge_x = x + (step_w - badge_h) / 2.0   # centre badge over box

        # Number badge
        add_shape_text(slide, MSO_SHAPE.OVAL, badge_x, y_box, badge_h, badge_h,
                       num, colour, WHITE, font_size=13, bold=True)

        # Action box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y_box + badge_h + 0.08), Inches(step_w), Inches(step_h))
        box.fill.solid()
        box.fill.fore_color.rgb = DK2 if colour == DK2 else colour
        box.line.fill.background()
        # Drop shadow on action box via XML
        sp_pr = box._element.spPr
        ef_lst = etree.SubElement(sp_pr, qn("a:effectLst"))
        outer  = etree.SubElement(ef_lst, qn("a:outerShdw"),
                                  blurRad="63500", dist="38100", dir="2700000",
                                  algn="tl", rotWithShape="0")
        srgb = etree.SubElement(outer, qn("a:srgbClr"), val="000000")
        etree.SubElement(srgb, qn("a:alpha"), val="45000")

        tf = box.text_frame; tf.word_wrap = True
        tf.auto_size = None
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Pt(6); tf.margin_right = Pt(6)
        tf.margin_top  = Pt(4); tf.margin_bottom = Pt(4)
        p = tf.paragraphs[0]; p.text = action
        p.font.name = "Arial"; p.font.size = Pt(10); p.font.bold = False
        p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER

        # Impact label — wider than box so "SPG_NO_RESULTSET" fits on one line
        label_w = 1.90   # wider than step_w to prevent wrapping
        label_x = x - (label_w - step_w) / 2.0   # centre under action box
        imp_box = slide.shapes.add_textbox(
            Inches(label_x), Inches(y_label), Inches(label_w), Inches(0.55))
        imp_box.text_frame.word_wrap = False   # single line per \n segment
        for i, line in enumerate(impact.split('\n')):
            if i == 0:
                para = imp_box.text_frame.paragraphs[0]
            else:
                para = imp_box.text_frame.add_paragraph()
            para.text = line
            para.font.name = "Arial"; para.font.size = Pt(10); para.font.bold = True
            para.font.color.rgb = PASS_GREEN; para.alignment = PP_ALIGN.CENTER

        # Connector arrow — prominent SF_BLUE filled RIGHT_ARROW with shadow
        if num != "5":
            arrow_x = x + step_w + arrow_gap
            arrow_y = y_box + badge_h + 0.08 + (step_h - arrow_h) / 2.0
            arr = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                Inches(arrow_x), Inches(arrow_y), Inches(arrow_w), Inches(arrow_h))
            arr.fill.solid()
            arr.fill.fore_color.rgb = SF_BLUE
            arr.line.fill.background()
            # Drop shadow on arrow
            arr_sp = arr._element.spPr
            arr_ef = etree.SubElement(arr_sp, qn("a:effectLst"))
            arr_os = etree.SubElement(arr_ef, qn("a:outerShdw"),
                                      blurRad="38100", dist="25400", dir="2700000",
                                      algn="tl", rotWithShape="0")
            arr_c  = etree.SubElement(arr_os, qn("a:srgbClr"), val="000000")
            etree.SubElement(arr_c, qn("a:alpha"), val="40000")
            arr.text_frame.paragraphs[0].text = ""   # no text inside arrow

        x += step_w + arrow_gap * 2 + arrow_w

    verify_slide(slide, prs, 99)


def slide_thank_you(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[28])
    set_ph(slide, 1, "THANK\nYOU")
    verify_slide(slide, prs, 100)


# ── Schema Audit Slides ───────────────────────────────────────────────────────

def fetch_schema_audit_data(spg_host, spg_user, spg_password, spg_db, schema_run_number):
    """Fetch schema-only audit results for the given run number."""
    try:
        import psycopg2, psycopg2.extras
        conn = psycopg2.connect(
            host=spg_host, port=5432, user=spg_user, password=spg_password,
            dbname=spg_db, sslmode='require', connect_timeout=15)
        cur = conn.cursor(cursor_factory=psycopg2.extras.RealDictCursor)
        cur.execute('''
            SELECT object_type, source_schema, test_verdict, COUNT(*) AS cnt
            FROM validation.validation_result
            WHERE run_number = %s
            GROUP BY object_type, source_schema, test_verdict
            ORDER BY object_type, source_schema, test_verdict
        ''', (schema_run_number,))
        rows = [dict(r) for r in cur.fetchall()]

        cur.execute('''
            SELECT total_objects, pass_count, fail_count, skip_count, run_status,
                   run_started_at
            FROM validation.validation_run WHERE run_number = %s
        ''', (schema_run_number,))
        meta = dict(cur.fetchone() or {})
        conn.close()
        return rows, meta
    except Exception as e:
        print(f"  WARN: Could not fetch schema audit data: {e}")
        return [], {}


def _pct(p, f):
    return f"{p / (p + f) * 100:.0f}%" if (p + f) else "N/A"


def slide_schema_audit_kpi(prs, meta, rows):
    """Slide: Schema Audit KPIs — 4 stat tiles."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    PASS_V    = {'PASS', 'PASS_DML_PROC', 'PASS_WRITE_PROC', 'WRITE_EXPECTED_FAIL'}
    FAIL_V    = {'FAIL', 'SPG_ERROR', 'SPG_NO_RESULTSET', 'MSSQL_ERROR',
                 'BOTH_FAILED', 'WRITE_SPG_ERROR', 'WRITE_BOTH_FAILED', 'WRITE_MSSQL_ERROR', 'ERROR', 'WARN'}
    MISSING_V = {'MSSQL_ONLY', 'SPG_ONLY'}

    total  = meta.get('total_objects', 0)
    passed = meta.get('pass_count', 0)
    failed = meta.get('fail_count', 0)
    skip   = meta.get('skip_count', 0)
    rate   = _pct(passed, failed)

    # Count PROC_TO_FUNC
    ptf_pass = sum(r['cnt'] for r in rows
                   if r['object_type'] == 'PROC_TO_FUNC'
                   and r['test_verdict'] in PASS_V)

    slide   = prs.slides.add_slide(prs.slide_layouts[0])
    slide_w = prs.slide_width.inches
    slide_h = prs.slide_height.inches

    set_ph(slide, 0, "Part 1: Schema Audit")
    set_ph(slide, 1, "Object existence, column names, parameter parity — no data execution")

    SNOW_BLUE = RGBColor(0x29, 0xB5, 0xE8)
    GREEN     = RGBColor(0x00, 0xA6, 0x5A)

    stats = [
        (str(total),  "Objects Audited",       SNOW_BLUE),
        (rate,        "Pass Rate",              GREEN),
        (str(passed), "Passed",                 GREEN),
        (str(ptf_pass), "PROC→FUNC (migrated)", SNOW_BLUE),
    ]

    n      = len(stats)
    card_w = 2.8
    gap    = (slide_w - n * card_w) / (n + 1)
    top    = 2.1
    ht     = 2.0

    for i, (val, lbl, col) in enumerate(stats):
        x = gap + i * (card_w + gap)
        add_shape_text(slide, 1, x, top, card_w, ht,
                       text=val, font_size=40, bold=True, color=col,
                       v_anchor='middle', h_align='center')
        add_shape_text(slide, 1, x, top + 1.05, card_w, 0.65,
                       text=lbl, font_size=14, bold=False,
                       color=RGBColor(0x44, 0x44, 0x44),
                       v_anchor='top', h_align='center')

    verify_slide(slide, prs, 99)


def slide_schema_audit_table(prs, rows):
    """Slide: Schema Audit — breakdown by object type × schema."""
    from pptx.util import Inches, Pt
    from collections import defaultdict

    PASS_V    = {'PASS', 'PASS_DML_PROC', 'PASS_WRITE_PROC', 'WRITE_EXPECTED_FAIL'}
    FAIL_V    = {'FAIL', 'SPG_ERROR', 'SPG_NO_RESULTSET', 'MSSQL_ERROR',
                 'BOTH_FAILED', 'WRITE_SPG_ERROR', 'WRITE_BOTH_FAILED', 'WRITE_MSSQL_ERROR', 'ERROR', 'WARN'}
    MISSING_V = {'MSSQL_ONLY', 'SPG_ONLY', 'SKIPPED'}

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_ph(slide, 0, "Schema Audit: Results by Object Type")
    set_ph(slide, 1, "Column parity · parameter parity · existence · PROC_TO_FUNC")

    # Aggregate: object_type → {pass, fail, missing}
    agg = defaultdict(lambda: {'pass': 0, 'fail': 0, 'skip': 0})
    for r in rows:
        ot = r['object_type'].upper()
        v  = r['test_verdict']
        cnt = r['cnt']
        if v in PASS_V:
            agg[ot]['pass'] += cnt
        elif v in FAIL_V:
            agg[ot]['fail'] += cnt
        else:
            agg[ot]['skip'] += cnt

    # Sort by type
    TYPE_ORDER = ['TABLE', 'VIEW', 'PROCEDURE', 'FUNCTION', 'PROC_TO_FUNC', 'TRIGGER']
    ordered = [(t, agg[t]) for t in TYPE_ORDER if t in agg]
    for t, g in sorted(agg.items()):
        if t not in TYPE_ORDER:
            ordered.append((t, g))

    headers = ["Object Type", "Total", "Pass", "Fail/Error", "Skip/Missing", "Pass %"]
    data_rows = []
    for otype, g in ordered:
        p, f, s = g['pass'], g['fail'], g['skip']
        data_rows.append([
            otype.replace('PROC_TO_FUNC', 'PROC→FUNC'),
            str(p + f + s), str(p), str(f), str(s), _pct(p, f)
        ])
    # Totals row
    tp = sum(g['pass'] for _, g in ordered)
    tf = sum(g['fail'] for _, g in ordered)
    ts = sum(g['skip'] for _, g in ordered)
    data_rows.append(["TOTAL", str(tp+tf+ts), str(tp), str(tf), str(ts), _pct(tp, tf)])

    add_table_style2(slide, 0.6, 1.4, 12.0, min(0.35 * (len(data_rows) + 1), 4.5),
                     headers, data_rows)
    verify_slide(slide, prs, 9)


# ── Markdown report ───────────────────────────────────────────────────────────

def generate_markdown(client_name, author, run_date,
                      runs, schema_summary, failures, verdict_cats,
                      kpis, mssql_counts, spg_catalog_counts,
                      run_nums, output_path, mssql_db=""):
    """Generate a Markdown version of the migration validation report."""

    lines = []
    a = lines.append  # shorthand

    # ── Header ────────────────────────────────────────────────────────────────
    a(f"# Migration Validation Report — {client_name}")
    a(f"")
    a(f"**Author:** {author}  ")
    a(f"**Date:** {run_date}  ")
    a(f"**Source:** MSSQL {mssql_db}  ")
    a(f"**Target:** Snowflake Postgres  ")
    a(f"**Validation Runs:** {', '.join(str(r) for r in run_nums)}  ")
    a(f"")

    # ── Migration at a Glance ─────────────────────────────────────────────────
    a(f"---")
    a(f"")
    a(f"## Migration at a Glance")
    a(f"")
    a(f"| KPI | Value |")
    a(f"|-----|-------|")
    a(f"| Total MSSQL Objects | {kpis['total_mssql_objects']} |")
    a(f"| Total SPG Objects | {kpis['spg_objects']} |")
    a(f"| Pass Rate | **{kpis['pass_rate']}%** |")
    a(f"| Schemas Tested | {kpis['schemas']} ({', '.join(kpis['schemas_list'])}) |")
    a(f"| Total PASS | {kpis['total_pass']} |")
    a(f"| Total FAIL | {kpis['total_fail']} |")
    a(f"")

    # ── Object count by schema + type ─────────────────────────────────────────
    a(f"---")
    a(f"")
    a(f"## Object Count by Schema and Type")
    a(f"")
    a(f"| Schema | Type | # MSSQL | # SPG | Passed | Failed | Pass % |")
    a(f"|--------|------|---------|-------|--------|--------|--------|")

    seen = set()
    all_types = ['TABLE', 'VIEW', 'PROCEDURE', 'FUNCTION', 'TRIGGER']
    schemas = sorted({r['source_schema'] for r in schema_summary
                      if r.get('source_schema') not in ('validation', None)})

    # From validation results
    for row in schema_summary:
        schema = row.get('source_schema') or ''
        otype  = row.get('object_type') or ''
        if schema == 'validation':
            continue
        key = (schema, otype)
        seen.add(key)
        total   = row.get('total') or 0
        passed  = row.get('passed') or 0
        failed  = row.get('failed') or 0
        ms_cnt  = mssql_counts.get(key) or mssql_counts.get((schema.lower(), otype)) or 'n/a'
        spg_cnt = spg_catalog_counts.get(key) or spg_catalog_counts.get((schema.lower(), otype)) or 'n/a'
        pct     = f"{round(passed / total * 100)}%" if total else 'n/a'
        a(f"| {schema} | {otype} | {ms_cnt} | {spg_cnt} | {passed} | {failed} | {pct} |")

    # Catalog-only rows (TABLE/VIEW not in validation results — use structural match)
    for (schema, otype), cnt in sorted(mssql_counts.items()):
        if (schema, otype) not in seen and schema not in ('sys', 'information_schema'):
            spg_cnt_raw = spg_catalog_counts.get((schema, otype)) or spg_catalog_counts.get((schema.lower(), otype))
            spg_cnt = spg_cnt_raw if spg_cnt_raw is not None else 0
            spg_display = str(spg_cnt) if spg_cnt_raw is not None else 'n/a'
            # Structural pass/fail: objects present in both = passed, missing = failed
            passed = min(cnt, spg_cnt) if spg_cnt_raw is not None else 0
            failed = max(cnt - spg_cnt, 0) if spg_cnt_raw is not None else cnt
            pct = f"{round(passed/cnt*100)}%" if cnt > 0 and spg_cnt_raw is not None else 'n/a'
            a(f"| {schema} | {otype} | {cnt} | {spg_display} | {passed} | {failed} | {pct} |")
    a(f"")

    # ── Validation Runs (moved after schema/type table) ───────────────────────
    a(f"---")
    a(f"")
    a(f"### Validation Runs")
    a(f"")
    a(f"| Run # | Started | Object Type | Objects | Pass | Fail | Skip | Status |")
    a(f"|------:|---------|-------------|--------:|-----:|-----:|-----:|--------|")
    for r in runs:
        started = str(r.get('run_started_at', ''))[:10]
        note = r.get('notes', '')
        obj_type = ('TRIGGER' if 'Trigger' in note
                    else 'VIEW' if 'View' in note
                    else 'PROCEDURE/FUNCTION' if 'Procedure' in note or 'function' in note or 'output comparison' in note.lower()
                    else 'WRITE PROCEDURES' if 'rollback' in note.lower() or 'write' in note.lower()
                    else 'MIXED')
        status = r.get('run_status') or (
            'ALL_PASS' if (r.get('fail_count') or 0) == 0 and (r.get('error_count') or 0) == 0
            else 'PARTIAL_PASS' if (r.get('pass_count') or 0) > 0 else 'FAIL')
        a(f"| {r['run_number']} | {started} | {obj_type} | "
          f"{r.get('total_objects') or 0} | {r.get('pass_count') or 0} | "
          f"{r.get('fail_count') or 0} | {r.get('skip_count') or 0} | {status} |")
    a(f"")

    # ── Failure categories ────────────────────────────────────────────────────
    a(f"---")
    a(f"")
    a(f"## Failure Categories")
    a(f"")
    VERDICT_LABELS = {
        'FAIL':           'Data / column mismatch',
        'SPG_ERROR':      'SPG execution error',
        'SPG_NO_RESULTSET': 'SPG proc cannot return result set',
        'MSSQL_ERROR':    'MSSQL execution error',
        'BOTH_FAILED':    'Both sides failed',
        'MSSQL_ONLY':     'In MSSQL but not in SPG',
        'SPG_ONLY':       'In SPG but not in MSSQL',
    }
    a(f"| Verdict | Count | Description |")
    a(f"|---------|-------|-------------|")
    for v in verdict_cats:
        label = VERDICT_LABELS.get(v['test_verdict'], v['test_verdict'])
        a(f"| {v['test_verdict']} | {v['cnt']} | {label} |")
    a(f"")

    # ── Remediation priorities ────────────────────────────────────────────────
    a(f"---")
    a(f"")
    a(f"## Remediation Priorities")
    a(f"")

    FIXES = [
        ('BOOL_INT_MISMATCH',
         'BIT columns compared as boolean',
         'Replace `WHERE col = true/false` with `WHERE col = 1/0` in migrated procs/views'),
        ('MISSING_COLUMN',
         'Column reference not found in SPG',
         'Audit LLM-mutated column names (e.g. `mg.majorid`, `fdl.new`) against MSSQL source'),
        ('PROC_NO_RESULTSET',
         'PG PROCEDURE cannot return rows',
         'Convert `CREATE PROCEDURE` to `CREATE FUNCTION … RETURNS TABLE(…)` for read-only procs'),
        ('VIEW_COLUMN_RENAME',
         'SPG view column names diverge from MSSQL',
         'Fix `allergentag` → `allergentags` and similar pluralisation mutations in `_POST_LLM_PATCHES_VIEWS`'),
        ('STG_JOB_INFRASTRUCTURE',
         'stg.* procs need active MicrosLoad job',
         'Populate `stg.MicrosLoadStatus` with a test job record before validating stg export procs'),
    ]
    for i, (code, title, action) in enumerate(FIXES, 1):
        a(f"### {i}. {title}")
        a(f"")
        a(f"**Code:** `{code}`  ")
        a(f"**Fix:** {action}")
        a(f"")

    # ── Failed object details ─────────────────────────────────────────────────
    a(f"---")
    a(f"")
    a(f"## Appendix — Failed Object Details")
    a(f"")

    # Group failures by schema
    by_schema: dict = {}
    for f in failures:
        s = f.get('source_schema') or 'unknown'
        by_schema.setdefault(s, []).append(f)

    if not by_schema:
        a(f"_No failures recorded in the selected runs._")
        a(f"")
    else:
        for schema in sorted(by_schema.keys()):
            a(f"### Schema: `{schema}`")
            a(f"")
            a(f"| Object | Type | Verdict | Issue |")
            a(f"|--------|------|---------|-------|")
            for row in sorted(by_schema[schema], key=lambda x: (x.get('object_type',''), x.get('object_name',''))):
                obj   = row.get('object_name') or ''
                otype = row.get('object_type') or ''
                verd  = row.get('test_verdict') or ''
                issue = (row.get('top_issue') or '').replace('\n', ' ').replace('|', '\\|')[:120]
                a(f"| `{obj}` | {otype} | {verd} | {issue} |")
            a(f"")

    # ── Footer ────────────────────────────────────────────────────────────────
    a(f"---")
    a(f"")
    a(f"_Generated by Cortex Code Migration Validator  |  {author}  |  {run_date}_")
    a(f"")

    content = '\n'.join(lines)
    with open(output_path, 'w', encoding='utf-8') as fh:
        fh.write(content)
    print(f"Saved: {output_path}")
    return output_path


# ── Main ──────────────────────────────────────────────────────────────────────

def build_deck(client_name, author, run_date,
               spg_host, spg_user, spg_password, spg_db="postgres",
               mssql_host="localhost", mssql_port=1433, mssql_user="SA",
               mssql_password="", mssql_db="",
               run_numbers=None, output_path=None, schema_run=0):

    print(f"\nFetching validation data from SPG: {spg_host[:55]}...")
    runs, schema_summary, failures, verdict_cats, run_nums = fetch_validation_data(
        spg_host, spg_user, spg_password, spg_db, run_numbers)

    print(f"Fetching MSSQL object counts from: {mssql_host}:{mssql_port}/{mssql_db}...")
    mssql_counts = fetch_mssql_counts(mssql_host, mssql_port, mssql_user,
                                      mssql_password, mssql_db) if mssql_db else {}

    print(f"Fetching SPG catalog counts...")
    spg_catalog_counts = fetch_spg_counts(spg_host, spg_user, spg_password, spg_db)

    kpis = compute_kpis(runs, schema_summary, mssql_counts, spg_catalog_counts)

    print(f"Runs: {run_nums}  |  Pass rate: {kpis['pass_rate']}%  |  Schemas: {kpis['schemas_list']}")
    print(f"Building deck for: {client_name}  |  Author: {author}  |  Date: {run_date}\n")

    prs = Presentation(TEMPLATE)
    while len(prs.slides) > 0:
        sldId = prs.slides._sldIdLst[0]
        rId = (sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
               or sldId.get('r:id'))
        if rId: prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(sldId)

    # Build slides
    slide_cover(prs, client_name, run_date, author)
    slide_testing_approach(prs, client_name=client_name, run_date=run_date, mssql_db=mssql_db)
    slide_methodology(prs)

    # ── Part 1: Schema Validation (DDL Structure) ─────────────────────────────
    slide_section_divider(
        prs, "PART 1", "Schema Validation (DDL Structure)",
        "Checks whether every MSSQL object was deployed to Snowflake Postgres")
    slide_part1_schema_deployment(prs, mssql_counts, spg_catalog_counts)
    slide_part1_by_schema(prs, mssql_counts, spg_catalog_counts)

    # ── Part 2: Behavioral Validation (Live Execution) ────────────────────────
    slide_section_divider(
        prs, "PART 2", "Behavioral Validation (Live Execution)",
        "Executes objects on both sides and compares outputs, row counts and behavior")
    slide_exec_summary(prs, kpis)
    slide_schema_table(prs, schema_summary, mssql_counts, spg_catalog_counts)
    slide_pass_rate_visual(prs, schema_summary, mssql_counts, spg_catalog_counts)
    slide_failure_categories(prs, verdict_cats)
    slide_remediation(prs, verdict_cats)

    # Appendix divider before failure detail tables
    app_slide = prs.slides.add_slide(prs.slide_layouts[18])
    set_ph(app_slide, 1, "APPENDIX\nFAILED OBJECT DETAILS")

    sn = 9
    # api failures
    sn = slide_failure_detail(prs, failures, ['api'], sn)
    # dbo + stg failures
    sn = slide_failure_detail(prs, failures, ['dbo', 'stg'], sn)

    slide_thank_you(prs)

    verify_deck(prs)

    # Save — honour REPORT_OUTPUT env var when set by run.py
    if output_path is None:
        output_path = os.environ.get('REPORT_OUTPUT')
    if output_path is None:
        gdrive = os.path.expanduser("~/Google Drive/My Drive")
        safe_client = client_name.replace(' ', '_').replace('/', '-')
        fname = f"Migration_Validation_{run_date.replace('-','')}.pptx"
        if os.path.isdir(gdrive):
            output_path = os.path.join(gdrive, fname)
        else:
            output_path = os.path.expanduser(f"~/Downloads/{fname}")

    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"\nSaved: {output_path}")
    if "Google Drive" in output_path:
        print("→ Will sync to Google Drive. Open via drive.google.com")

    # ── Markdown report (generate alongside .pptx if not already present) ────
    md_path = os.path.splitext(output_path)[0] + ".md"
    if not os.path.exists(md_path):
        generate_markdown(
            client_name=client_name,
            author=author,
            run_date=run_date,
            runs=runs,
            schema_summary=schema_summary,
            failures=failures,
            verdict_cats=verdict_cats,
            kpis=kpis,
            mssql_counts=mssql_counts,
            spg_catalog_counts=spg_catalog_counts,
            run_nums=run_nums,
            output_path=md_path,
            mssql_db=mssql_db,
        )
    else:
        print(f"Skipped (already exists): {md_path}")

    return output_path


# ── CLI entry point ───────────────────────────────────────────────────────────

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Migration Validation Report Generator")
    parser.add_argument("--client",       default=os.environ.get("CLIENT_NAME", "Client"))
    parser.add_argument("--author",       default=os.environ.get("AUTHOR", "Rekha Khandhadia"))
    parser.add_argument("--date",         default=datetime.date.today().strftime("%Y%m%d"))
    parser.add_argument("--spg-host",     default=os.environ.get("SPG_HOST", ""))
    parser.add_argument("--spg-user",     default=os.environ.get("SPG_USER", "snowflake_admin"))
    parser.add_argument("--spg-password", default=os.environ.get("SPG_PASSWORD", ""))
    parser.add_argument("--spg-db",         default=os.environ.get("SPG_DATABASE", "postgres"))
    parser.add_argument("--mssql-host",     default=os.environ.get("MSSQL_HOST", "localhost"))
    parser.add_argument("--mssql-port",     default=os.environ.get("MSSQL_PORT", "1433"))
    parser.add_argument("--mssql-user",     default=os.environ.get("MSSQL_USER", "SA"))
    parser.add_argument("--mssql-password", default=os.environ.get("MSSQL_PASSWORD", ""))
    parser.add_argument("--mssql-db",       default=os.environ.get("MSSQL_DATABASE", ""))
    parser.add_argument("--run-numbers",    default=None,
                        help="Comma-separated run numbers, e.g. 1,2,3 (default: latest 3)")
    parser.add_argument("--schema-run",     type=int, default=0,
                        help="Schema-only audit run number for Part 1 slides (0 = omit)")
    parser.add_argument("--output",         default=None,
                        help="Output path for .pptx (default: Google Drive or ~/Downloads)")
    args = parser.parse_args()

    run_nums = [int(x) for x in args.run_numbers.split(",")] if args.run_numbers else None

    build_deck(
        client_name=args.client,
        author=args.author,
        run_date=args.date,
        spg_host=args.spg_host,
        spg_user=args.spg_user,
        spg_password=args.spg_password,
        spg_db=args.spg_db,
        mssql_host=args.mssql_host,
        mssql_port=int(args.mssql_port),
        mssql_user=args.mssql_user,
        mssql_password=args.mssql_password,
        mssql_db=args.mssql_db,
        run_numbers=run_nums,
        output_path=args.output,
        schema_run=args.schema_run,
    )
